import os
import gc
import random
import tempfile
import zipfile
import uuid
import numpy as np
import torch
from PIL import Image
from diffusers import QwenImageLayeredPipeline
import gradio as gr

# --- 依赖库检查 ---
# 如果运行报错，请在终端执行：pip install python-pptx psd-tools
try:
    from pptx import Presentation
    from psd_tools import PSDImage
except ImportError:
    print("请先安装依赖: pip install python-pptx psd-tools")
    exit()

# ================= 配置与显存优化 (关键修改部分) =================
# 1. 设置环境变量
os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"
os.environ["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True"
os.environ["HF_HOME"] = "/root/autodl-fs/.cache/huggingface"

print("正在清理残留内存...")
torch.cuda.empty_cache()
gc.collect()

print("正在加载模型 (优化模式)...")
# 2. 优化加载方式：直接使用 bfloat16 和 low_cpu_mem_usage
pipeline = QwenImageLayeredPipeline.from_pretrained(
    "Qwen/Qwen-Image-Layered",
    torch_dtype=torch.bfloat16,
    low_cpu_mem_usage=True,
)

# 3. 启用智能 CPU Offload (替代 .to("cuda"))
# 这样模型会根据需要在 CPU 和 GPU 间切换，节省大量显存
pipeline.enable_model_cpu_offload()

# 4. 尝试开启 VAE 优化 (能开则开，不能开跳过)
try:
    if hasattr(pipeline, "vae"):
        pipeline.vae.enable_slicing()
        pipeline.vae.enable_tiling()
        print("VAE 显存优化已开启")
except Exception:
    pass

pipeline.set_progress_bar_config(disable=None)
MAX_SEED = np.iinfo(np.int32).max
# =============================================================

def imagelist_to_pptx(img_files):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        emu = inch * 914400
        return int(emu)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(img_path, left, top, width=px_to_emu(img_width_px), height=px_to_emu(img_height_px))

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name

def imagelist_to_psd(img_files):
    layers = []
    for path in img_files:
        layers.append(Image.open(path).convert('RGBA'))

    width, height = layers[0].size
    psd = PSDImage.new(mode='RGBA', size=(width, height))

    for i, img in enumerate(layers):
        name = f"Layer {i + 1}"
        layer = psd.create_pixel_layer(image=img, name=name)
        psd.append(layer)

    with tempfile.NamedTemporaryFile(suffix=".psd", delete=False) as tmp:
        psd.save(tmp.name)
        return tmp.name

def infer(input_image,
          seed=777,
          randomize_seed=False,
          prompt=None,
          neg_prompt=" ",
          true_guidance_scale=4.0,
          num_inference_steps=50,
          layer=4,
          cfg_norm=True,
          use_en_prompt=True):

    # 每次推理前清理显存
    gc.collect()
    torch.cuda.empty_cache()

    if randomize_seed:
        seed = random.randint(0, MAX_SEED)

    if isinstance(input_image, list):
        input_image = input_image[0]

    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError("Unsupported input_image type: %s" % type(input_image))

    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device='cpu').manual_seed(seed), # Generator 放 CPU 更安全
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,

        # 【关键修改】从 640 改为 448，防止最后一步 Killed
        "resolution": 512,

        "cfg_normalize": cfg_norm,
        "use_en_prompt": use_en_prompt,
    }
    print(f"Start Inference with Resolution: {inputs['resolution']}")

    try:
        with torch.inference_mode():
            output = pipeline(**inputs)
            output_images = output.images[0]
    except Exception as e:
        raise gr.Error(f"推理失败: {e} (可能是显存不足，请降低分辨率或步数)")

    output = []
    temp_files = []
    for i, image in enumerate(output_images):
        output.append(image)
        # Save to temp file for export
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        image.save(tmp.name)
        temp_files.append(tmp.name)

    # Generate PPTX
    pptx_path = imagelist_to_pptx(temp_files)

    # Generate ZIP
    with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as tmp:
        with zipfile.ZipFile(tmp.name, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, img_path in enumerate(temp_files):
                zipf.write(img_path, f"layer_{i+1}.png")
        zip_path = tmp.name

    # Generate PSD
    psd_path = imagelist_to_psd(temp_files)
    return output, pptx_path, zip_path, psd_path

# --- Gradio UI 部分保持不变 ---
# (为了节省篇幅，这里假设你使用同样的 Examples 列表，你可以把你的 examples 列表放在这里)
examples = [] # 这里填入你的图片路径，确保路径存在

with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.HTML('''<p align="center"><h1>Qwen-Image-Layered (AutoDL Optimized)</h1><p>''')
        gr.Markdown("""
                    The text prompt is intended to describe the overall content.
                    **Note:** Resolution has been optimized to 448 to prevent memory crashes.
                    """)
        with gr.Row():
            with gr.Column(scale=1):
                input_image = gr.Image(label="Input Image", image_mode="RGBA", type="pil") # 指定 type="pil" 更稳定

                with gr.Accordion("Advanced Settings", open=True): # 默认展开方便调试
                    prompt = gr.Textbox(label="Prompt (Optional)", lines=3)
                    neg_prompt = gr.Textbox(label="Negative Prompt", value=" ", lines=1)

                    seed = gr.Slider(label="Seed", minimum=0, maximum=MAX_SEED, step=1, value=0)
                    randomize_seed = gr.Checkbox(label="Randomize seed", value=True)

                    true_guidance_scale = gr.Slider(label="True guidance scale", minimum=1.0, maximum=10.0, value=4.0)
                    num_inference_steps = gr.Slider(label="Steps", minimum=1, maximum=50, value=20) # 默认先设20跑通再说
                    layer = gr.Slider(label="Layers", minimum=2, maximum=10, step=1, value=4)

                    cfg_norm = gr.Checkbox(label="Enable CFG normalization", value=True)
                    use_en_prompt = gr.Checkbox(label="Auto EN Prompt", value=True)

                run_button = gr.Button("Decompose!", variant="primary")

            with gr.Column(scale=2):
                gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
                with gr.Row():
                    export_file = gr.File(label="Download PPTX")
                    export_zip_file = gr.File(label="Download ZIP")
                    export_psd_file = gr.File(label="Download PSD")

    # 注意：这里去掉了 Examples 的 run_on_click=True，防止一打开就自动跑崩
    # 如果你有图片，可以把 examples 变量填进去
    if examples:
        gr.Examples(examples=examples, inputs=[input_image], outputs=[gallery, export_file, export_zip_file, export_psd_file], fn=infer)

    run_button.click(
        fn=infer,
        inputs=[
            input_image, seed, randomize_seed, prompt, neg_prompt,
            true_guidance_scale, num_inference_steps, layer, cfg_norm, use_en_prompt
        ],
        outputs=[gallery, export_file, export_zip_file, export_psd_file],
    )

print("ready to launch... http://localhost:6006")

demo.launch(
    server_name="0.0.0.0",
    server_port=6006,
    show_error=True # 显示详细报错
)
